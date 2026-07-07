#============================================
#NAME: test_remediation_accuracy_checkpoints
#============================================

#===========================
#DIRECTIONS:
#===========================

"""
Process and validate document conversions by verifying page count and preview quality,
confirming complete page conversion across output formats, checking AI-generated image
descriptions, evaluating OCR and formatting accuracy against the source document,
and validating table tagging and internal link preservation where applicable
"""
# ==========================
# IMPORTS
# ==========================
from __future__ import annotations

from itertools import count
import difflib
import re
import unicodedata
import zipfile
from pathlib import Path
from urllib.parse import urlparse

from odf.opendocument import load
from odf.table import Table, TableRow, TableCell
from odf.teletype import extractText

from striprtf.striprtf import rtf_to_text
import docx
import pandas as pd
from bs4 import BeautifulSoup
import openpyxl
from pptx import Presentation
import xlrd
from pypdf import PdfReader
from pypdf.generic import (
    ArrayObject,
    DictionaryObject,
    IndirectObject,
    NameObject,
    TextStringObject,
    ContentStream,
)
from playwright.sync_api import TimeoutError as PlaywrightTimeoutError
import pytest
# ==========================
# CONSTANTS
# ==========================
WEBSITE = "https://test.scribeit.io/"

BASE_DIR = Path(__file__).resolve().parent

TEST_DOCS_DIR = BASE_DIR / "test_documents" / "various_test_documents_for_projects"

DEBUG_DIR = BASE_DIR / "debug"
DOWNLOADS_DIR = BASE_DIR / "downloads"

# AI Description Settings
AI_DESC_MARKERS = [
    "AI-generated description",
    "AI generated description",
    "image description",
]


# Preview Rules
PREVIEW_MAX_PAGES = 3

# Force all conversions to a specific output format
FORCE_OUTPUT_FORMAT: str | None = None
# ============================================================
# SUPPORTED FORMATS
# ============================================================

FORMAT_MAP = {
    ".pdf": "PDF",
    ".doc": "Word",
    ".docx": "Word",
    ".rtf": "Word",
    ".ppt": "PDF",
    ".pptx": "PDF",
    ".txt": "Word",
    ".odt": "Word",
    ".ods": "Word",
    ".xls": "Word",
    ".xlsx": "Word",
    ".brf": "Braille",
    ".html": "Webpage",
    ".htm": "Webpage",
    ".xhtml": "Webpage",
    ".epub": "EPUB",
    ".daisy": "DAISY",
    ".mobi": "MOBI",
}

SUPPORTED_IMG = {
    ".bmp": "BMP",
    ".gif": "GIF",
    ".jpg": "JPG",
    ".jpeg": "JPG",
    ".png": "PNG",
    ".tiff": "TIFF",
    ".tif": "TIFF",
}

FORMAT_ALIASES = {
    "Word": ["Word", "DOCX", "Microsoft Word", "Word (DOCX)", "Microsoft Word (DOCX)"],
    "PDF": ["PDF", "Tagged PDF", "Accessible PDF"],
}

SUPPORTED_EXT = tuple(FORMAT_MAP.keys()) + tuple(SUPPORTED_IMG.keys())
# ==========================
# LOGGING + DEBUG DUMPS
# ==========================
def format_label_candidates(label: str) -> list[str]:
    return FORMAT_ALIASES.get(label, [label])


def _safe_write(path: Path, data: str) -> None:
    try:
        path.write_text(data, encoding="utf-8", errors="ignore")
    except Exception:
        pass

def dump_ui_debug(page, label: str) -> None:
    try:
        DEBUG_DIR.mkdir(parents=True, exist_ok=True)
    except Exception:
        return

    try:
        page.screenshot(path=str(DEBUG_DIR / f"{label}.png"), full_page=True)
    except Exception:
        pass

    try:
        _safe_write(DEBUG_DIR / f"{label}.html", page.content())
    except Exception:
        pass

    try:
        lines = []
        for i, fr in enumerate(page.frames):
            try:
                lines.append(f"frame[{i}] name={fr.name!r} url={fr.url!r}")
            except Exception:
                lines.append(f"frame[{i}] <unable to read url/name>")
        _safe_write(DEBUG_DIR / f"{label}_frames.txt", "\n".join(lines))
    except Exception:
        pass

    try:
        for i, fr in enumerate(page.frames[:5]):
            try:
                if not fr.url or fr.url.startswith("about:"):
                    continue
                _safe_write(DEBUG_DIR / f"{label}_frame_{i}.html", fr.content())
            except Exception:
                continue
    except Exception:
        pass

    print(f"[DEBUG] Saved debug artifacts for '{label}' in ./debug/", flush=True)
# ==========================
# SCOPES (page + frames)
# ==========================
def scopes(page):
    scopes_ = [page]
    try:
        scopes_.extend(list(page.frames))
    except Exception:
        pass
    return scopes_


def first_clickable(locator, timeout_ms: int = 5000) -> bool:
    try:
        if locator is None:
            return False
        if locator.count() <= 0:
            return False
        locator.first.wait_for(state="visible", timeout=timeout_ms)
        locator.first.click(timeout=timeout_ms)
        return True
    except Exception:
        return False
# ==========================
# OUTPUT FORMAT PICKER
# ==========================
def decide_output_format_label(source_path: Path) -> str:
    if FORCE_OUTPUT_FORMAT:
        return FORCE_OUTPUT_FORMAT

    ext = source_path.suffix.lower()
    if ext in FORMAT_MAP:
        return FORMAT_MAP[ext]
    if ext in SUPPORTED_IMG:
        return "PDF"  # default output for images
    raise ValueError(f"Unsupported extension: {ext}")
# ==========================
# UI FLOW HELPERS
# ==========================
def goto_home(page) -> None:
    try:
        page.goto(WEBSITE, wait_until="domcontentloaded", timeout=60_000)
    except Exception:
        dump_ui_debug(page, "goto_home_failed")
        raise

def upload_document_any_scope(page, file_path: str, timeout_ms: int = 30_000) -> None:
    last_err = None
    for sc in scopes(page):
        try:
            loc = sc.locator("input[type='file']")
            if loc.count() > 0:
                loc.first.wait_for(state="attached", timeout=timeout_ms)
                loc.first.set_input_files(file_path, timeout=timeout_ms)
                return
        except Exception as e:
            last_err = e
            continue
    dump_ui_debug(page, "upload_input_not_found")
    raise RuntimeError(f"Could not find file input to upload. Last error: {last_err}")

def click_start_any_scope(page, timeout_ms: int = 60_000) -> None:
    start_pat = re.compile(r"^\s*start\s*$", re.I)

    for sc in scopes(page):
        try:
            btn = sc.get_by_role("button", name=start_pat)
            if first_clickable(btn, timeout_ms=timeout_ms):
                return
        except Exception:
            pass

    for sc in scopes(page):
        try:
            loc = sc.get_by_text(start_pat, exact=False)
            if first_clickable(loc, timeout_ms=timeout_ms):
                return
        except Exception:
            pass

    dump_ui_debug(page, "start_not_found")
    raise RuntimeError("Could not find/click Start button.")

def wait_for_after_start(page, timeout_ms: int = 120_000) -> None:
    try:
        page.wait_for_url(re.compile(r".*/documents/.*", re.I), timeout=timeout_ms)
        try:
            page.wait_for_load_state("domcontentloaded", timeout=30_000)
        except Exception:
            pass
        try:
            page.wait_for_load_state("networkidle", timeout=30_000)
        except Exception:
            pass
        return
    except Exception:
        pass

    try:
        for sc in scopes(page):
            sc.get_by_text(re.compile(r"This document has", re.I)).first.wait_for(timeout=30_000)
            return
    except Exception:
        pass

    convert_like = re.compile(r"^\s*(convert|go|remediate|run)\b", re.I)
    for sc in scopes(page):
        try:
            btn = sc.get_by_role("button", name=convert_like)
            btn.first.wait_for(state="visible", timeout=30_000)
            return
        except Exception:
            continue

    dump_ui_debug(page, "after_start_unknown_state")
    raise RuntimeError("After Start, did not detect doc route, preview text, or convert-like button.")

def ensure_on_document_main_view(page, timeout_ms: int = 60_000) -> None:
    try:
        url = page.url or ""
    except Exception:
        return

    m = re.match(
        r"^(https://test\.scribeit\.io/documents/[^/]+)(/.*)?$",
        url
    )
    if not m:
        return

    base = m.group(1)
    if "/html_stream" in url:
        page.goto(base, wait_until="domcontentloaded", timeout=timeout_ms)
        try:
            page.wait_for_load_state("networkidle", timeout=15_000)
        except Exception:
            pass

def wait_for_preview_processed_any_scope(page, timeout_ms: int = 180_000) -> None:
    last_err = None
    for sc in scopes(page):
        try:
            sc.get_by_text(re.compile(r"This document has", re.I)).first.wait_for(timeout=timeout_ms)
            return
        except Exception as e:
            last_err = e
            continue
    dump_ui_debug(page, "preview_processed_not_found")
    raise RuntimeError(f"Preview completion text not found. Last error: {last_err}")

def open_select_format_required(page, timeout_ms: int = 30_000) -> None:
    """
    REQUIRED: click the 'Select format' / 'Choose format' control
    AND verify the dropdown/menu is actually open (visible).
    """
    open_pats = [
        re.compile(r"select\s*format", re.I),
        re.compile(r"choose\s*format", re.I),
        re.compile(r"output\s*format", re.I),
    ]

    start = page.evaluate("Date.now()")
    last_err = None

    def menu_is_open() -> bool:
        for sc in scopes(page):
            try:
                menu = sc.locator("[role='listbox'], [role='menu']")
                if menu.count() > 0:
                    try:
                        menu.first.wait_for(state="visible", timeout=500)
                        return True
                    except Exception:
                        pass
            except Exception:
                pass
        return False

    while True:
        if menu_is_open():
            return

        for sc in scopes(page):
            for role in ("button", "combobox"):
                for pat in open_pats:
                    try:
                        loc = sc.get_by_role(role, name=pat)
                        if first_clickable(loc, timeout_ms=1500):
                            if menu_is_open():
                                return
                    except Exception as e:
                        last_err = e

            for pat in open_pats:
                try:
                    label = sc.get_by_text(pat, exact=False)
                    if label.count() > 0:
                        container = label.first.locator(
                            "xpath=ancestor-or-self::*[self::button or @role='button' or @role='combobox' or self::div][1]"
                        )
                        if first_clickable(container, timeout_ms=1500) and menu_is_open():
                            return
                        if first_clickable(label, timeout_ms=1500) and menu_is_open():
                            return
                except Exception as e:
                    last_err = e

        if page.evaluate("Date.now()") - start > timeout_ms:
            dump_ui_debug(page, "select_format_control_not_found")
            raise RuntimeError(f"Could not open 'Select format' dropdown. Last error: {last_err}")

        page.wait_for_timeout(250)

def click_format_option_required(page, format_label: str, timeout_ms: int = 30_000) -> None:
    wanted_contains = re.compile(rf"\b{re.escape(format_label)}\b", re.I)

    start = page.evaluate("Date.now()")
    last_err = None

    while True:
        for sc in scopes(page):
            try:
                sel = sc.locator("select")
                if sel.count() > 0:
                    sel.first.select_option(label=format_label, timeout=2000)
                    return
            except Exception as e:
                last_err = e

            try:
                menu = sc.locator("[role='listbox'], [role='menu']")
                if menu.count() > 0:
                    for role in ("option", "menuitem", "menuitemradio"):
                        try:
                            opt = sc.get_by_role(role, name=wanted_contains)
                            if first_clickable(opt, timeout_ms=2000):
                                return
                        except Exception as e:
                            last_err = e

                    try:
                        opt2 = menu.locator("*").filter(has_text=wanted_contains)
                        if first_clickable(opt2, timeout_ms=2000):
                            return
                    except Exception as e:
                        last_err = e
            except Exception as e:
                last_err = e

        if page.evaluate("Date.now()") - start > timeout_ms:
            dump_ui_debug(page, "format_option_not_found")
            raise RuntimeError(f"Format option '{format_label}' not found/clickable. Last error: {last_err}")

        page.wait_for_timeout(250)

def click_convert_like_any_scope(page, timeout_ms: int = 60_000) -> None:
    patterns = [
        re.compile(r"^\s*convert\s*$", re.I),
        re.compile(r"^\s*go\s*$", re.I),
        re.compile(r"^\s*remediate\b", re.I),
        re.compile(r"\bremediate\s+full\s+document\b", re.I),
        re.compile(r"\bfull\s+document\b", re.I),
        re.compile(r"\brun\b", re.I),
        re.compile(r"\bprocess\b", re.I),
    ]

    for pat in patterns:
        for sc in scopes(page):
            try:
                btn = sc.get_by_role("button", name=pat)
                if first_clickable(btn, timeout_ms=timeout_ms):
                    return
            except Exception:
                pass

    for pat in patterns:
        for sc in scopes(page):
            try:
                loc = sc.get_by_text(pat, exact=False)
                if first_clickable(loc, timeout_ms=timeout_ms):
                    return
            except Exception:
                pass

    dump_ui_debug(page, "convert_like_not_found")
    raise RuntimeError("Could not find a Convert/Go/Remediate trigger.")

def click_download_if_present_any_scope(page, timeout_ms: int = 180_000) -> bool:
    pat = re.compile(r"\bdownload\b", re.I)
    deadline = page.evaluate("Date.now()") + timeout_ms

    while page.evaluate("Date.now()") < deadline:
        for sc in scopes(page):
            for role in ("button", "link"):
                try:
                    loc = sc.get_by_role(role, name=pat)
                    if first_clickable(loc, timeout_ms=1500):
                        return True
                except Exception:
                    pass

        for sc in scopes(page):
            try:
                loc = sc.get_by_text(pat, exact=False)
                if loc.count() > 0:
                    container = loc.first.locator(
                        "xpath=ancestor-or-self::*[self::button or self::a or @role='button' or @role='link'][1]"
                    )
                    if first_clickable(container, timeout_ms=1500) or first_clickable(loc, timeout_ms=1500):
                        return True
            except Exception:
                pass

        page.wait_for_timeout(500)

    return False

def wait_until_format_ready(page, timeout_ms: int = 180_000) -> None:
    # wait until preview processing appears finished
    wait_for_preview_processed_any_scope(page, timeout_ms=timeout_ms)

    # make sure you are on the main document screen
    ensure_on_document_main_view(page)

    try:
        # use your existing helper that already searches more broadly
        open_select_format_required(page, timeout_ms=timeout_ms)
        return

    except Exception:
        # optional: save debug artifacts for this specific failure
        dump_ui_debug(page, "format_control_not_ready")

        raise RuntimeError(
            "Preview appeared but Select Format button did not become enabled."
        )
# ==========================
# CHECKLIST HELPERS
# ==========================
def checklist(step_num: int, label: str, *, ok: bool, details: str | None = None) -> None:
    status = "✅" if ok else "❌"
    msg = f"[STEP {step_num}] {status} {label}"
    if details:
        msg += f" — {details}"
    print(msg, flush=True)


def run_checklist_step(step_num: int, label: str, fn) -> None:
    try:
        fn()
        checklist(step_num, label, ok=True)
    except Exception as e:
        checklist(step_num, label, ok=False, details=str(e))
        raise

def run_flow_and_download(page, source_path: Path, downloads_dir: Path) -> tuple[str, Path | None]:
    """
    Pipeline:
      Choose file -> Start -> Open Select format -> Pick format -> Convert -> Download
    """
    downloads_dir.mkdir(parents=True, exist_ok=True)
    out_label = decide_output_format_label(source_path)

    run_checklist_step(1, "Open Scribe home", lambda: goto_home(page))
    run_checklist_step(2, "Document uploaded", lambda: upload_document_any_scope(page, str(source_path)))
    run_checklist_step(3, "Start clicked", lambda: click_start_any_scope(page))
    run_checklist_step(4, "Preview loaded", lambda: wait_for_after_start(page))
    run_checklist_step(5, "Format controls ready",
                       lambda: (ensure_on_document_main_view(page), wait_until_format_ready(page)))

    # Step 6: Pick output format
    picked = False
    last_err: Exception | None = None

    for _attempt in range(2):  # retry once
        try:
            open_select_format_required(page)
        except Exception as e:
            last_err = e
            continue

        for candidate in format_label_candidates(out_label):
            try:
                click_format_option_required(page, candidate)
                picked = True
                break
            except Exception as e:
                last_err = e
                continue

        if picked:
            break

    if not picked:
        checklist(6, "Output format selected", ok=False,
                  details=f"Could not select: {out_label}. Last error: {last_err}")
        dump_ui_debug(page, "choose_format_failed_all_candidates")
        raise RuntimeError(f"Could not select any format option for: {out_label}")

    checklist(6, "Output format selected", ok=True)

    # Step 7: Convert / Remediate -> Download
    converted_path: Path | None = None
    before_local = {p.name for p in downloads_dir.glob("*")}

    sys_downloads_dir = Path.home() / "Downloads"
    before_sys = {p.name for p in sys_downloads_dir.glob("*")} if sys_downloads_dir.is_dir() else set()

    # Attempt A: Convert triggers immediate download
    try:
        with page.expect_download(timeout=30_000) as dl_info:
            click_convert_like_any_scope(page, timeout_ms=60_000)
        dl = dl_info.value
        target = downloads_dir / dl.suggested_filename
        dl.save_as(str(target))
        converted_path = target
    except PlaywrightTimeoutError:
        pass
    except Exception:
        pass

    # Attempt B: wait for explicit Download button
    if converted_path is None:
        try:
            with page.expect_download(timeout=180_000) as dl_info:
                if not click_download_if_present_any_scope(page, timeout_ms=180_000):
                    raise RuntimeError("Download button did not appear after Convert/Remediate.")
            dl = dl_info.value
            target = downloads_dir / dl.suggested_filename
            dl.save_as(str(target))
            converted_path = target
        except PlaywrightTimeoutError:
            pass
        except Exception:
            pass

    # Attempt C: filesystem fallback
    if converted_path is None:
        for _ in range(180):  # ~180s
            page.wait_for_timeout(1000)

            new_local = [
                p for p in downloads_dir.glob("*")
                if p.is_file() and p.name not in before_local and not p.name.endswith(".crdownload")
            ]
            if new_local:
                converted_path = max(new_local, key=lambda p: p.stat().st_mtime)
                break

            if sys_downloads_dir.is_dir():
                new_sys = [
                    p for p in sys_downloads_dir.glob("*")
                    if p.is_file() and p.name not in before_sys and not p.name.endswith(".crdownload")
                ]
                if new_sys:
                    latest = max(new_sys, key=lambda p: p.stat().st_mtime)
                    target = downloads_dir / latest.name
                    try:
                        target.write_bytes(latest.read_bytes())
                        converted_path = target
                    except Exception:
                        converted_path = latest
                    break

    if converted_path is not None:
        checklist(7, "Convert pressed and file downloaded", ok=True)
    else:
        checklist(7, "Convert pressed and file downloaded", ok=False, details="No download detected.")

    return out_label, converted_path
# =============================================================================
# LOW-LEVEL PDF STRUCTURE HELPERS
# =============================================================================
def _resolve(obj):
    return obj.get_object() if isinstance(obj, IndirectObject) else obj

def _name_to_str(n):
    n = _resolve(n)
    if isinstance(n, NameObject):
        return str(n).lstrip("/")
    if isinstance(n, str) and n.startswith("/"):
        return n.lstrip("/")
    return None

def _iter_struct_elems_in_order(struct_tree_root: DictionaryObject):
    def walk_k(k):
        k = _resolve(k)
        if isinstance(k, DictionaryObject):
            yield k
            kk = k.get("/K")
            if kk is not None:
                yield from walk_k(kk)
        elif isinstance(k, (ArrayObject, list)):
            for item in k:
                yield from walk_k(item)

    k0 = struct_tree_root.get("/K")
    if k0 is not None:
        yield from walk_k(k0)

def _heading_level(struct_elem: DictionaryObject):
    s = _name_to_str(struct_elem.get("/S"))
    if not s:
        return None
    if len(s) == 2 and s[0] == "H" and s[1].isdigit():
        lvl = int(s[1])
        if 1 <= lvl <= 6:
            return lvl
    return None

def _safe_text_from_tags(struct_elem: DictionaryObject) -> str | None:
    def norm_text(t):
        t = _resolve(t)
        if isinstance(t, (TextStringObject, str)):
            s = str(t).strip()
            return s if s else None
        return None

    return norm_text(struct_elem.get("/ActualText")) or norm_text(struct_elem.get("/Alt"))

def _page_ref_key(obj):
    obj = _resolve(obj)
    if isinstance(obj, IndirectObject):
        return (obj.idnum, obj.generation)
    ref = getattr(obj, "indirect_reference", None)
    if isinstance(ref, IndirectObject):
        return (ref.idnum, ref.generation)
    return None

def _build_page_index_map(reader: PdfReader) -> dict[tuple[int, int], int]:
    out: dict[tuple[int, int], int] = {}
    for i, pg in enumerate(reader.pages):
        k = _page_ref_key(pg)
        if k:
            out[k] = i
    return out

def _extract_text_from_tj_operand(val) -> str:
    val = _resolve(val)
    if isinstance(val, (TextStringObject, str)):
        return str(val)
    return ""

def _extract_text_from_TJ_array(arr) -> str:
    arr = _resolve(arr)
    if not isinstance(arr, (ArrayObject, list)):
        return ""
    parts: list[str] = []
    for it in arr:
        it = _resolve(it)
        if isinstance(it, (TextStringObject, str)):
            parts.append(str(it))
    return "".join(parts)

def _mcid_text_map_for_page(reader: PdfReader, page) -> dict[int, str]:
    out: dict[int, list[str]] = {}
    try:
        contents = page.get_contents()
        if contents is None:
            return {}
        cs = ContentStream(contents, reader)
    except Exception:
        return {}

    mcid_stack: list[int | None] = []

    def current_mcid() -> int | None:
        for v in reversed(mcid_stack):
            if isinstance(v, int):
                return v
        return None

    for operands, operator in getattr(cs, "operations", []):
        op = operator

        if op == b"BDC":
            mcid_val = None
            props = None
            for o in operands:
                o2 = _resolve(o)
                if isinstance(o2, DictionaryObject):
                    props = o2
                    break
            if isinstance(props, DictionaryObject):
                mc = _resolve(props.get("/MCID"))
                if isinstance(mc, int):
                    mcid_val = mc
            mcid_stack.append(mcid_val)
            continue

        if op == b"BMC":
            mcid_stack.append(None)
            continue

        if op == b"EMC":
            if mcid_stack:
                mcid_stack.pop()
            continue

        mcid = current_mcid()
        if mcid is None:
            continue

        try:
            if op == b"Tj" or op == b"'":
                s = _extract_text_from_tj_operand(operands[0]) if operands else ""
            elif op == b'"':
                s = _extract_text_from_tj_operand(operands[-1]) if operands else ""
            elif op == b"TJ":
                s = _extract_text_from_TJ_array(operands[0]) if operands else ""
            else:
                s = ""
        except Exception:
            s = ""

        if s:
            out.setdefault(mcid, []).append(s)

    return {k: "".join(v).strip() for k, v in out.items() if "".join(v).strip()}

def _mcid_lookup(reader: PdfReader, mcid: int, pg_hint, page_index_map, page_mcid_maps) -> str | None:
    if pg_hint is not None:
        key = _page_ref_key(pg_hint)
        if key and key in page_index_map:
            idx = page_index_map[key]
            text = page_mcid_maps.get(idx, {}).get(mcid)
            if text:
                return text

    for idx, mp in page_mcid_maps.items():
        if mcid in mp and mp[mcid]:
            return mp[mcid]
    return None

def _struct_k_to_text_fragments(reader: PdfReader, k, pg_hint, page_index_map, page_mcid_maps) -> list[str]:
    frags: list[str] = []
    k = _resolve(k)

    if isinstance(k, int):
        frag = _mcid_lookup(reader, k, pg_hint, page_index_map, page_mcid_maps)
        if frag:
            frags.append(frag)
        return frags

    if isinstance(k, (ArrayObject, list)):
        for it in k:
            frags.extend(_struct_k_to_text_fragments(reader, it, pg_hint, page_index_map, page_mcid_maps))
        return frags

    if isinstance(k, DictionaryObject):
        mc = _resolve(k.get("/MCID"))
        if isinstance(mc, int):
            frag = _mcid_lookup(reader, mc, _resolve(k.get("/Pg")) or pg_hint, page_index_map, page_mcid_maps)
            if frag:
                frags.append(frag)
            return frags

        if _name_to_str(k.get("/Type")) == "OBJR" or k.get("/OBJ") is not None:
            return frags

        kk = k.get("/K")
        if kk is not None:
            frags.extend(
                _struct_k_to_text_fragments(
                    reader,
                    kk,
                    _resolve(k.get("/Pg")) or pg_hint,
                    page_index_map,
                    page_mcid_maps,
                )
            )
        return frags

    return frags

def extract_tagged_text_reading_order(pdf_path: Path) -> str:
    try:
        reader = PdfReader(str(pdf_path))
        root = _resolve(reader.trailer.get("/Root"))
        struct_tree_root = _resolve(root.get("/StructTreeRoot")) if isinstance(root, DictionaryObject) else None
        if not isinstance(struct_tree_root, DictionaryObject):
            return ""

        page_index_map = _build_page_index_map(reader)
        page_mcid_maps: dict[int, dict[int, str]] = {}
        for i, pg in enumerate(reader.pages):
            page_mcid_maps[i] = _mcid_text_map_for_page(reader, pg)

        out_parts: list[str] = []
        last = ""

        for elem in _iter_struct_elems_in_order(struct_tree_root):
            if not isinstance(elem, DictionaryObject):
                continue

            t = _safe_text_from_tags(elem)
            if not t:
                k = elem.get("/K")
                pg_hint = _resolve(elem.get("/Pg"))
                if k is not None:
                    frags = _struct_k_to_text_fragments(reader, k, pg_hint, page_index_map, page_mcid_maps)
                    t = " ".join([f for f in frags if f]).strip()

            if not t:
                continue

            t = " ".join(str(t).split()).strip()
            if not t or t == last:
                continue

            role = _name_to_str(elem.get("/S")) or ""
            if role.startswith("H") or role in {"P", "LI", "LBody", "TD", "TH", "Caption", "Table", "TR"}:
                out_parts.append(t)
            else:
                if out_parts:
                    out_parts[-1] = (out_parts[-1] + " " + t).strip()
                else:
                    out_parts.append(t)

            last = t

        return "\n".join(out_parts).strip()
    except Exception:
        return ""
# =============================================================================
# STRING NORMALIZATION + SIMILARITY
# =============================================================================
def normalize_unicode(text: str) -> str:
    return unicodedata.normalize("NFKC", text)

def _norm(s: str) -> str:
    s = (s or "").strip()
    s = s.replace("\n", " ").replace("\t", " ").replace("\r", " ")
    s = " ".join(s.split())
    s = normalize_unicode(s).lower()
    return s

def sequence_similarity(a: list[str], b: list[str]) -> float:
    a2 = [_norm(x) for x in (a or []) if _norm(x)]
    b2 = [_norm(x) for x in (b or []) if _norm(x)]
    return difflib.SequenceMatcher(None, a2, b2).ratio()

def text_similarity(a: str, b: str) -> float:
    return difflib.SequenceMatcher(None, _norm(a), _norm(b)).ratio()
# =============================================================================
# OCR / TEXT QUALITY
# =============================================================================
def levenshtein_distance(a: str, b: str) -> int:
    if a == b:
        return 0
    if not a:
        return len(b)
    if not b:
        return len(a)

    prev = list(range(len(b) + 1))
    for i, ca in enumerate(a, start=1):
        cur = [i]
        for j, cb in enumerate(b, start=1):
            ins = cur[j - 1] + 1
            dele = prev[j] + 1
            sub = prev[j - 1] + (0 if ca == cb else 1)
            cur.append(min(ins, dele, sub))
        prev = cur
    return prev[-1]

def calculate_cer(reference_text: str, ocr_text: str) -> float:
    if not reference_text or not ocr_text:
        return 1.0
    return levenshtein_distance(reference_text, ocr_text) / max(1, len(reference_text))

def extract_text_from_file(path: Path) -> str:
    sfx = path.suffix.lower()

    if sfx in (".txt", ".log"):
        return path.read_text(encoding="utf-8", errors="ignore")

    if sfx in (".html", ".htm", ".xhtml"):
        html = path.read_text(encoding="utf-8", errors="ignore")
        soup = BeautifulSoup(html, "html.parser")
        return soup.get_text(" ", strip=True)

    if sfx == ".docx":
        doc = docx.Document(str(path))
        parts = []
        for p in doc.paragraphs:
            t = (p.text or "").strip()
            if t:
                parts.append(t)
        return "\n".join(parts)

    if sfx == ".pdf":
        try:
            tagged = extract_tagged_text_reading_order(path)
            if tagged:
                return tagged
            reader = PdfReader(str(path))
            return "\n".join((pg.extract_text() or "") for pg in reader.pages)
        except Exception:
            return ""

    if sfx == ".rtf":
        raw_rtf = path.read_text(encoding="utf-8", errors="ignore")
        return rtf_to_text(raw_rtf)

    if sfx == ".xls":
        tables = extract_all_tables(xls_path=path)
        return "\n".join(
            " | ".join(str(cell) for cell in row)
            for table in tables
            for row in table.values.tolist()
        )

    if sfx == ".pptx":
        prs = Presentation(str(path))
        parts = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_value = shape.text.strip()
                    if text_value:
                        parts.append(text_value)

        return "\n".join(parts)

    return ""

def _sample_text(s: str, max_chars: int = 4000) -> str:
    s = s or ""
    if len(s) <= max_chars * 2:
        return s
    return s[:max_chars] + "\n...\n" + s[-max_chars:]

def compare_ocr_quality_best_effort(source_path: Path, converted_path: Path) -> dict:
    source_text = extract_text_from_file(source_path)
    converted_text = extract_text_from_file(converted_path)

    if not source_text and not converted_text:
        return {"status": "UNKNOWN", "cer": None, "wer": None, "notes": "Insufficient text to compare."}
    if source_text and not converted_text:
        return {"status": "FAIL", "cer": None, "wer": None, "notes": "Converted text empty; extraction/OCR degraded."}

    MAX_FULL_COMPARE_CHARS = 2500

    if len(source_text) > MAX_FULL_COMPARE_CHARS or len(converted_text) > MAX_FULL_COMPARE_CHARS:
        src_s = _sample_text(source_text, max_chars=4000)
        cnv_s = _sample_text(converted_text, max_chars=4000)

        char_ratio = difflib.SequenceMatcher(None, _norm(src_s), _norm(cnv_s)).ratio()
        cer = 1.0 - char_ratio

        src_words = _norm(src_s).split()
        cnv_words = _norm(cnv_s).split()
        word_ratio = difflib.SequenceMatcher(None, src_words, cnv_words).ratio()
        wer = 1.0 - word_ratio

        return {
            "status": "PASS",
            "cer": cer,
            "wer": wer,
            "notes": "Large document detected; CER/WER approximated using sampled text (fast mode).",
        }

    cer = calculate_cer(source_text, converted_text)

    src_words = _norm(source_text).split()
    cnv_words = _norm(converted_text).split()
    word_ratio = difflib.SequenceMatcher(None, src_words, cnv_words).ratio()
    wer = 1.0 - word_ratio

    return {"status": "PASS", "cer": cer, "wer": wer, "notes": "Compared extracted text (not raster OCR)."}
# =============================================================================
# PAGE COUNT HELPERS
# =============================================================================
def _docx_pages_from_app_xml(docx_path: Path) -> int | None:
    """
    Best-effort DOCX page count from docProps/app.xml.
    Exists if the DOCX was saved by Word (or a tool that writes it).
    """
    try:
        with zipfile.ZipFile(str(docx_path), "r") as zf:
            xml = zf.read("docProps/app.xml").decode("utf-8", errors="ignore")
        m = re.search(r"<Pages>(\d+)</Pages>", xml)
        return int(m.group(1)) if m else None
    except Exception:
        return None

def file_page_count_best_effort(path: Path) -> int | None:
    """
    Returns a real file-based page count when that concept exists.
    Reflowable formats (EPUB/MOBI/DAISY) -> None (use UI fallback).
    """
    sfx = path.suffix.lower()

    if sfx == ".pdf":
        try:
            return len(PdfReader(str(path)).pages)
        except Exception:
            return None

    if sfx == ".docx":
        return _docx_pages_from_app_xml(path)

    return None

def page_count_from_ui_any_scope(page) -> int:
    wait_for_preview_processed_any_scope(page)
    for sc in scopes(page):
        try:
            text = sc.get_by_text(re.compile(r"This document has", re.I)).first.text_content()
            match = re.search(r"\d+", text or "")
            if match:
                return int(match.group(0))
        except Exception:
            continue
    raise ValueError("Could not parse page count from UI text.")

def preview_pages_visible_count_any_scope(page) -> int | None:
    """
    Best-effort: count preview pages shown in the UI.
    Falls back to parsing UI text if needed.
    """
    try:
        wait_for_preview_processed_any_scope(page, timeout_ms=180_000)
    except Exception:
        return None

    page_selectors = [
        "[data-page]",
        "[data-page-number]",
        "[data-page-index]",
        "[data-testid*='page']",
        ".react-pdf__Page",
        "canvas[data-page-number]",
        "canvas[data-page-index]",
        "canvas[aria-label*='Page' i]",
        "[aria-label^='Page ' i]",
        "img[alt^='Page ' i]",
    ]

    preview_containers = [
        "[data-testid*='preview']",
        "[class*='preview']",
        "[class*='Preview']",
        "[role='document']",
        "main",
    ]

    def plausible(n: int) -> bool:
        return 1 <= n <= 500

    for sc in scopes(page):
        for cont_sel in preview_containers:
            try:
                cont = sc.locator(cont_sel)
                if cont.count() <= 0:
                    continue
                cont0 = cont.first

                for sel in page_selectors:
                    loc = cont0.locator(sel)
                    n = loc.count()
                    if n <= 0:
                        continue
                    try:
                        loc.first.wait_for(state="attached", timeout=2_500)
                    except Exception:
                        continue
                    n = loc.count()
                    if plausible(n):
                        return n
            except Exception:
                continue

        for sel in page_selectors:
            try:
                loc = sc.locator(sel)
                n = loc.count()
                if n <= 0:
                    continue
                try:
                    loc.first.wait_for(state="attached", timeout=2_500)
                except Exception:
                    continue
                n = loc.count()
                if plausible(n):
                    return n
            except Exception:
                continue

    for sc in scopes(page):
        try:
            t = sc.get_by_text(re.compile(r"This document has\s+\d+\s+pages", re.I)).first.text_content()
            m = re.search(r"(\d+)", t or "")
            if m:
                return int(m.group(1))
        except Exception:
            continue

    return None

def preview_page_flag(preview_pages: int | None, max_preview_pages: int = PREVIEW_MAX_PAGES) -> str | None:
    if preview_pages is None:
        return "Preview page count could not be detected (UI selector needs adjustment)."
    if preview_pages > max_preview_pages:
        return f"Preview exceeds expected size: {preview_pages} pages shown (max {max_preview_pages})."
    return None
# =============================================================================
# AI IMAGE DESCRIPTION CHECKS
# =============================================================================
def count_ai_markers_in_text(text: str) -> int:
    return sum(
        len(re.findall(re.escape(marker), text or "", flags=re.IGNORECASE))
        for marker in AI_DESC_MARKERS
    )

def count_unlabeled_images_html(html_text: str) -> int:
    soup = BeautifulSoup(html_text, "html.parser")
    return sum(1 for img in soup.find_all("img") if not (img.get("alt") or "").strip())

def count_unlabeled_figures_pdf(pdf_path: Path) -> int:
    try:
        reader = PdfReader(str(pdf_path))
        root = _resolve(reader.trailer.get("/Root"))
        struct_tree_root = _resolve(root.get("/StructTreeRoot")) if isinstance(root, DictionaryObject) else None
        if not isinstance(struct_tree_root, DictionaryObject):
            return 0

        unlabeled = 0
        for elem in _iter_struct_elems_in_order(struct_tree_root):
            if _name_to_str(elem.get("/S")) != "Figure":
                continue
            alt = _safe_text_from_tags(elem)
            if not alt:
                unlabeled += 1
        return unlabeled
    except Exception:
        return 0

def expected_unlabeled_graphics(source_path: Path) -> int:
    sfx = source_path.suffix.lower()
    if sfx in (".html", ".htm", ".xhtml"):
        return count_unlabeled_images_html(source_path.read_text(encoding="utf-8", errors="ignore"))
    if sfx == ".pdf":
        return count_unlabeled_figures_pdf(source_path)
    return 0
# =============================================================================
# LINKS + TABLE TAGGING CHECKS
# =============================================================================
def is_valid_url(url: str) -> bool:
    try:
        p = urlparse(url)
        return bool(p.scheme and p.netloc)
    except Exception:
        return False

def links_supported(path: Path) -> bool:
    return path.suffix.lower() in {".pdf", ".html", ".htm", ".xhtml"}

def extract_pdf_links(pdf_path: str | Path) -> set[str]:
    reader = PdfReader(str(pdf_path))
    links: set[str] = set()

    for page in reader.pages:
        annots = _resolve(page.get("/Annots"))
        if not annots:
            continue
        ann_list = annots if isinstance(annots, (ArrayObject, list)) else [annots]

        for a in ann_list:
            a = _resolve(a)
            if not isinstance(a, DictionaryObject):
                continue
            if _name_to_str(a.get("/Subtype")) != "Link":
                continue

            action = _resolve(a.get("/A"))
            if isinstance(action, DictionaryObject) and _name_to_str(action.get("/S")) == "URI":
                uri = _resolve(action.get("/URI"))
                if isinstance(uri, (TextStringObject, str)):
                    s = str(uri).strip()
                    if s:
                        links.add(s)

            dest = _resolve(a.get("/Dest"))
            if dest is not None:
                links.add(str(dest))

    return links

def extract_html_links(html_text: str) -> set[str]:
    soup = BeautifulSoup(html_text, "html.parser")
    return {(a.get("href") or "").strip() for a in soup.find_all("a") if (a.get("href") or "").strip()}

def extract_pdf_table_tag_metadata(pdf_path: str | Path) -> dict:
    reader = PdfReader(str(pdf_path))
    root = _resolve(reader.trailer.get("/Root"))
    struct_tree_root = _resolve(root.get("/StructTreeRoot")) if isinstance(root, DictionaryObject) else None
    if not isinstance(struct_tree_root, DictionaryObject):
        return {"tables": 0, "has_caption": False, "has_th": False}

    tables = 0
    has_caption = False
    has_th = False

    def walk(node):
        nonlocal has_caption, has_th, tables
        node = _resolve(node)
        if isinstance(node, DictionaryObject):
            role = _name_to_str(node.get("/S"))
            if role == "Table":
                tables += 1
            if role == "Caption":
                has_caption = True
            if role == "TH":
                has_th = True
            kk = node.get("/K")
            if kk is not None:
                walk(kk)
        elif isinstance(node, (ArrayObject, list)):
            for it in node:
                walk(it)

    walk(struct_tree_root.get("/K"))
    return {"tables": tables, "has_caption": has_caption, "has_th": has_th}
# =============================================================================
# FILE PICKING
# =============================================================================
def test_files() -> list[Path]:
    if not TEST_DOCS_DIR.is_dir():
        return []
    return sorted([p for p in TEST_DOCS_DIR.iterdir() if p.is_file() and p.name.lower().endswith(SUPPORTED_EXT)])
# =============================================================================
# EXTRACTION (HEADINGS / LISTS / TABLES)
# =============================================================================
def extract_all_headings(*, file_path: str | Path | None = None, html_text: str | None = None) -> list[str]:
    headings: list[str] = []

    if html_text is not None:
        soup = BeautifulSoup(html_text, "html.parser")
        return [h.get_text(strip=True) for h in soup.find_all(re.compile(r"^h[1-6]$")) if h.get_text(strip=True)]

    if file_path is None:
        return []

    path = Path(file_path)
    suffix = path.suffix.lower()

    if suffix == ".pdf":
        reader = PdfReader(str(path))
        root = _resolve(reader.trailer.get("/Root"))
        struct_tree_root = _resolve(root.get("/StructTreeRoot")) if isinstance(root, DictionaryObject) else None
        if not isinstance(struct_tree_root, DictionaryObject):
            return []
        for elem in _iter_struct_elems_in_order(struct_tree_root):
            if _heading_level(elem) is None:
                continue
            text = _safe_text_from_tags(elem)
            if text:
                headings.append(text)

    elif suffix in {".html", ".htm", ".xhtml"}:
        return extract_all_headings(html_text=path.read_text(encoding="utf-8", errors="ignore"))

    elif suffix in {".epub", ".zip"}:
        with zipfile.ZipFile(str(path), "r") as zf:
            for name in zf.namelist():
                if name.lower().endswith((".html", ".xhtml", ".htm")):
                    html = zf.read(name).decode("utf-8", errors="ignore")
                    headings.extend(extract_all_headings(html_text=html))

    elif suffix == ".docx":
        doc = docx.Document(str(path))
        for p in doc.paragraphs:
            t = (p.text or "").strip()
            if t and any(getattr(run, "bold", False) for run in p.runs):
                headings.append(t)

    return headings

def extract_all_list_items(*, pdf_path: str | Path | None = None, html_text: str | None = None,
                           zip_path: str | Path | None = None) -> list[str]:
    items: list[str] = []

    if html_text is not None:
        soup = BeautifulSoup(html_text, "html.parser")
        items = [li.get_text(" ", strip=True) for li in soup.find_all("li")]
        return items or ["No List Found"]

    if pdf_path is not None:
        reader = PdfReader(str(pdf_path))
        root = _resolve(reader.trailer.get("/Root"))
        struct_tree_root = _resolve(root.get("/StructTreeRoot")) if isinstance(root, DictionaryObject) else None
        if not isinstance(struct_tree_root, DictionaryObject):
            return ["No List Found"]
        for elem in _iter_struct_elems_in_order(struct_tree_root):
            if _name_to_str(elem.get("/S")) != "LI":
                continue
            text = _safe_text_from_tags(elem)
            if text:
                items.append(text.strip())
        return items or ["No List Found"]

    if zip_path is not None:
        with zipfile.ZipFile(str(zip_path), "r") as zf:
            for name in zf.namelist():
                if name.lower().endswith((".html", ".xhtml", ".htm")):
                    html = zf.read(name).decode("utf-8", errors="ignore")
                    soup = BeautifulSoup(html, "html.parser")
                    items.extend([li.get_text(" ", strip=True) for li in soup.find_all("li")])
        return items or ["No List Found"]

    return ["No List Found"]

def extract_all_tables (  *,
    html_text: str | None = None,
    zip_path: str | Path | None = None,
    pdf_path: str | Path | None = None,
    docx_path: str | Path | None = None,
    xlsx_path: str | Path | None = None,
    ods_path: str | Path | None = None,
    xls_path: str | Path | None = None,
    odt_path: str | Path | None = None

) -> list[pd.DataFrame]:
    tables: list[pd.DataFrame] = []

    if html_text is not None:
        soup = BeautifulSoup(html_text, "html.parser")
        for table_tag in soup.find_all("table"):
            headers = [th.get_text(strip=True) for th in table_tag.find_all("th")]
            rows = []
            for row in table_tag.find_all("tr"):
                cols = [td.get_text(strip=True) for td in row.find_all("td")]
                if cols:
                    rows.append(cols)
            if rows:
                tables.append(pd.DataFrame(rows, columns=headers if headers else None))
        return tables

    if zip_path is not None:
        with zipfile.ZipFile(str(zip_path), "r") as zf:
            for name in zf.namelist():
                if name.lower().endswith((".html", ".xhtml", ".htm")):
                    html = zf.read(name).decode("utf-8", errors="ignore")
                    tables.extend(extract_all_tables(html_text=html))
        return tables

    if pdf_path is not None:
        reader = PdfReader(str(pdf_path))
        root = _resolve(reader.trailer.get("/Root"))
        struct_tree_root = _resolve(root.get("/StructTreeRoot")) if isinstance(root, DictionaryObject) else None
        if not isinstance(struct_tree_root, DictionaryObject):
            return []

        for elem in _iter_struct_elems_in_order(struct_tree_root):
            if _name_to_str(elem.get("/S")) != "Table":
                continue

            rows = []
            k = _resolve(elem.get("/K"))
            kids = k if isinstance(k, (ArrayObject, list)) else ([k] if k is not None else [])
            for row in kids:
                row = _resolve(row)
                if not isinstance(row, DictionaryObject) or _name_to_str(row.get("/S")) != "TR":
                    continue

                row_cells = []
                row_kids = _resolve(row.get("/K"))
                row_children = row_kids if isinstance(row_kids, (ArrayObject, list)) else (
                    [row_kids] if row_kids else [])
                for cell in row_children:
                    cell = _resolve(cell)
                    if not isinstance(cell, DictionaryObject):
                        continue
                    if _name_to_str(cell.get("/S")) not in {"TH", "TD"}:
                        continue
                    cell_text = _safe_text_from_tags(cell) or ""
                    row_cells.append(cell_text.strip())

                if row_cells:
                    rows.append(row_cells)

            if rows:
                tables.append(pd.DataFrame(rows))

        return tables

    if docx_path is not None:
        doc = docx.Document(str(docx_path))
        for table in doc.tables:
            table_rows = []
            for row in table.rows:
                row_values = []
                for cell in row.cells:
                    row_values.append(" ".join((cell.text or "").strip().split()))
                table_rows.append(row_values)
            if table_rows:
                tables.append(pd.DataFrame(table_rows))
        return tables

    if ods_path is not None:
        doc = load(str(ods_path))

        for table in doc.getElementsByType(Table):
            rows = []

            for tr in table.getElementsByType(TableRow):
                row_data = []

                if ods_path is not None:
                    doc = load(str(ods_path))

                    for table in doc.getElementsByType(Table):
                        rows = []

                        for tr in table.getElementsByType(TableRow):
                            row_data = []

                            for tc in tr.getElementsByType(TableCell):
                                text_content = extractText(tc).strip()
                                row_data.append(text_content)

                            if any(row_data):
                                rows.append(row_data)

                        if rows:
                            tables.append(pd.DataFrame(rows))

                    return tables

                if any(row_data):
                    rows.append(row_data)

            if rows:
                tables.append(pd.DataFrame(rows))

        return tables

    if xlsx_path is not None:
        wb = openpyxl.load_workbook(str(xlsx_path), data_only=True)

        for sheet in wb.worksheets:
            rows = []

            for row in sheet.iter_rows(values_only=True):
                cleaned_row = []

                for cell in row:
                    if cell is None:
                        cleaned_row.append("")
                    else:
                        cleaned_row.append(str(cell))

                if any(cleaned_row):
                    rows.append(cleaned_row)

            if rows:
                tables.append(pd.DataFrame(rows))

        return tables

    if xls_path is not None:
        book = xlrd.open_workbook(str(xls_path))

        for sheet in book.sheets():
            rows = []

            for row_idx in range(sheet.nrows):
                row_values = []
                for col_idx in range(sheet.ncols):
                    cell_value = sheet.cell_value(row_idx, col_idx)
                    row_values.append("" if cell_value is None else str(cell_value))

                if any(v.strip() for v in row_values):
                    rows.append(row_values)

            if rows:
                tables.append(pd.DataFrame(rows))

        return tables

    return []
# =============================================================================
# NORMALIZATION + COMPARISON HELPERS
# =============================================================================
def normalize_all_headings(headings: list[str]) -> list[str]:
    out: list[str] = []
    for h in headings or []:
        hh = _norm(h)
        if hh and hh not in out:
            out.append(hh)
    return out

def normalize_list(items: list[str] | None) -> list[str]:
    out: list[str] = []
    for item in items or []:
        ii = _norm(item)
        if ii and ii not in out:
            out.append(ii)
    return out

def normalize_tables(tables: list[pd.DataFrame] | None) -> list[str]:
    normalized_rows: list[str] = []
    for df in tables or []:
        if df is None or df.empty:
            continue
        for _, row in df.iterrows():
            cleaned = [_norm("" if c is None else str(c)) for c in row]
            cleaned = [c for c in cleaned if c]
            if not cleaned:
                continue
            row_string = " | ".join(cleaned)
            if row_string not in normalized_rows:
                normalized_rows.append(row_string)
    return normalized_rows

def compare_documents(source_items: list[str], converted_items: list[str]) -> tuple[list[str], list[str]]:
    src = normalize_list(source_items)
    cnv_set = set(normalize_list(converted_items))
    matched, unmatched = [], []
    for item in src:
        (matched if item in cnv_set else unmatched).append(item)
    return matched, unmatched
# =============================================================================
# RUNNER HELPERS
# =============================================================================
def is_limited_text_comparison_file_type(path: Path) -> bool:
    return path.suffix.lower() in {
        ".bmp",
        ".gif",
        ".jpg",
        ".jpeg",
        ".png",
        ".tiff",
        ".tif",
        ".odt"
    }

def should_use_limited_text_comparison(source_path: Path, converted_path: Path) -> bool:
    return (
        is_limited_text_comparison_file_type(source_path)
        or is_limited_text_comparison_file_type(converted_path)
    )

def tables_for(p: Path) -> list[pd.DataFrame]:
    sfx = p.suffix.lower()

    if sfx == ".pdf":
        return extract_all_tables(pdf_path=p)

    if sfx in (".html", ".htm", ".xhtml"):
        return extract_all_tables(html_text=p.read_text(encoding="utf-8", errors="ignore"))

    if sfx in (".zip", ".epub"):
        return extract_all_tables(zip_path=p)

    if sfx == ".docx":
        return extract_all_tables(docx_path=p)

    if sfx == ".xlsx":
        return extract_all_tables(xlsx_path=p)

    if sfx == ".ods":
        return extract_all_tables(ods_path=p)

    if sfx == ".xls":
        return extract_all_tables(xls_path=p)

    if sfx == ".odt":
        return extract_all_tables(odt_path=p)
    return []

def lists_for(p: Path) -> list[str]:
    sfx = p.suffix.lower()
    if sfx == ".pdf":
        return extract_all_list_items(pdf_path=p)
    if sfx in (".html", ".htm", ".xhtml"):
        return extract_all_list_items(html_text=p.read_text(encoding="utf-8", errors="ignore"))
    if sfx in (".zip", ".epub"):
        return extract_all_list_items(zip_path=p)
    return ["No List Found"]

def fmt_pct_1dp(x: float | None) -> str:
    """Convert ratio (0..1) to percent with 1 decimal, e.g. 0.031 -> 3.1%"""
    if x is None:
        return "N/A"
    return f"{x * 100.0:.1f}%"
# =============================================================================
# CHECKPOINT RUNNER (ALL 7)
# =============================================================================
def run_checkpoints_for_file(context, source_path: Path) -> dict:
    result = {
        "source": str(source_path),
        "converted": None,
        "ui_page_count": None,
        "preview_pages_visible": None,
        "preview_flag": None,
        "pages_converted_ok": None,
        "source_file_pages": None,
        "converted_file_pages": None,
        "unlabeled_graphics_expected": None,
        "ai_desc_marker_count": None,
        "ai_desc_ok": None,
        "ocr_quality": None,
        "headings": None,
        "lists": None,
        "tables": None,
        "text_similarity": None,
        "table_tagging": None,
        "links": None,
        "error": None,
    }

    page = context.pages[0] if context.pages else context.new_page()

    ck = count(1)

    def ckprint(msg: str) -> None:
        print(f"[{next(ck)}] {msg}", flush=True)

    try:
        # FLOW + DOWNLOAD
        try:
            chosen_label, converted_path = run_flow_and_download(page, source_path, DOWNLOADS_DIR)
            print(f"[FLOW] Output format decided: {chosen_label}")
        except Exception as e:
            print(f"[FAIL] Flow failed: {e}")
            result["error"] = str(e)
            return result

        # (1) UI page count
        try:
            ui_pc = page_count_from_ui_any_scope(page)
            result["ui_page_count"] = ui_pc
            ckprint(f"UI PAGE COUNT (after preview): {ui_pc}")
        except Exception as e:
            print(f"[WARN] UI page count failed: {e}")

        # (2) Preview amount + flag
        pv_pages = preview_pages_visible_count_any_scope(page)
        result["preview_pages_visible"] = pv_pages
        flag = preview_page_flag(pv_pages, max_preview_pages=PREVIEW_MAX_PAGES)
        result["preview_flag"] = flag

        ckprint(f"PREVIEW PAGES SHOWN: {pv_pages}")
        ckprint(f"PREVIEW FLAG: {flag if flag else 'none'}")

        # Download result
        if not converted_path:
            ckprint("FAIL: No download detected after Convert-like action.")
            print("Check ./debug/ for artifacts.")
            return result

        result["converted"] = str(converted_path)
        ckprint(f"DOWNLOADED: {converted_path.name}")

        # (3) Confirm all pages converted
        ui_pages = result.get("ui_page_count")

        source_pages_file = file_page_count_best_effort(source_path)
        converted_pages_file = file_page_count_best_effort(converted_path)

        if source_path.suffix.lower() == ".docx" and ui_pages is not None:
            expected_pages = ui_pages
        else:
            expected_pages = source_pages_file if source_pages_file is not None else ui_pages

        actual_pages = converted_pages_file if converted_pages_file is not None else ui_pages

        result["source_file_pages"] = expected_pages
        result["converted_file_pages"] = actual_pages

        if expected_pages is None or actual_pages is None:
            result["pages_converted_ok"] = None
            ckprint("PAGES CONVERTED: UNKNOWN (missing file count + UI count)")
        else:
            ok = (actual_pages == expected_pages)
            result["pages_converted_ok"] = ok
            ckprint(f"PAGES CONVERTED: expected = {expected_pages}, actual = {actual_pages}")

        # (4) AI-generated image descriptions
        unlabeled_expected = expected_unlabeled_graphics(source_path)
        result["unlabeled_graphics_expected"] = unlabeled_expected

        converted_text = extract_text_from_file(converted_path)
        marker_count = count_ai_markers_in_text(converted_text)
        result["ai_desc_marker_count"] = marker_count

        ai_ok = (marker_count == unlabeled_expected)
        result["ai_desc_ok"] = ai_ok

        ckprint(f"UNLABELED GRAPHICS EXPECTED (source): {unlabeled_expected}")
        ckprint(f"AI DESCRIPTION MARKERS FOUND (converted): {marker_count}")
        ckprint(f"AI DESCRIPTIONS OK: {ai_ok}")

        # (5) OCR / text quality
        ocr_report = compare_ocr_quality_best_effort(source_path, converted_path)
        result["ocr_quality"] = ocr_report

        ckprint(
            f"OCR/TEXT QUALITY: {ocr_report.get('status')} | "
            f"Character Error Rate={fmt_pct_1dp(ocr_report.get('cer'))} | "
            f"Word Error Rate={fmt_pct_1dp(ocr_report.get('wer'))} | {ocr_report.get('notes')}"
        )

        # (6) Compare formatting (headings/lists/tables/text similarity)
        src_text = extract_text_from_file(source_path)
        cnv_text = extract_text_from_file(converted_path)

        limited_text_comparison = should_use_limited_text_comparison(source_path, converted_path)

        if limited_text_comparison:
            ttxt_sim = None
            text_warn = False
            result["text_similarity"] = {
                "similarity": ttxt_sim,
                "warn": text_warn,
                "limited": True,
                "note": "TEXT COMPARISON: limited / not reliable for this file type",
            }
            ckprint("OVERALL TEXT similarity=SKIPPED (limited / not reliable for this file type)")
            print("[INFO] TEXT COMPARISON: limited / not reliable for this file type", flush=True)
        else:
            ttxt_sim = text_similarity(src_text, cnv_text)
            result["text_similarity"] = {
                "similarity": ttxt_sim,
                "limited": False,
            }
            ckprint(f"OVERALL TEXT similarity={ttxt_sim:.1%}")


        src_h = normalize_all_headings(extract_all_headings(file_path=source_path))
        cnv_h = normalize_all_headings(extract_all_headings(file_path=converted_path))

        h_match, h_unmatch = compare_documents(src_h, cnv_h)
        h_sim = sequence_similarity(src_h, cnv_h)

        result["headings"] = {
            "similarity": h_sim,
            "matched": h_match,
            "unmatched": h_unmatch,
        }

        ckprint(
            f"HEADINGS similarity={h_sim:.1%} | "
            f"matched={len(h_match)} unmatched={len(h_unmatch)}"
        )

        src_l = normalize_list(lists_for(source_path))
        cnv_l = normalize_list(lists_for(converted_path))
        l_match, l_unmatch = compare_documents(src_l, cnv_l)
        l_sim = sequence_similarity(src_l, cnv_l)
        result["lists"] = {"similarity": l_sim, "matched": l_match, "unmatched": l_unmatch}
        ckprint(f"LISTS similarity={l_sim:.1%} | matched={len(l_match)} unmatched={len(l_unmatch)}")

        src_t = normalize_tables(tables_for(source_path))
        cnv_t = normalize_tables(tables_for(converted_path))
        t_match, t_unmatch = compare_documents(src_t, cnv_t)
        t_sim = sequence_similarity(src_t, cnv_t)
        result["tables"] = {"similarity": t_sim, "matched": t_match, "unmatched": t_unmatch}
        ckprint(f"TABLES similarity={t_sim:.1%} | matched={len(t_match)} unmatched={len(t_unmatch)}")


        # (7) Table tagging + links
        if source_path.suffix.lower() == ".pdf" and converted_path.suffix.lower() == ".pdf":
            src_meta = extract_pdf_table_tag_metadata(source_path)
            cnv_meta = extract_pdf_table_tag_metadata(converted_path)

            ok = True
            if src_meta["tables"] > 0 and cnv_meta["tables"] == 0:
                ok = False
            if src_meta["has_caption"] and not cnv_meta["has_caption"]:
                ok = False
            if src_meta["has_th"] and not cnv_meta["has_th"]:
                ok = False

            result["table_tagging"] = {"source": src_meta, "converted": cnv_meta, "ok": ok}
            ckprint(f"TABLE TAGGING (PDF) ok={ok} | {result['table_tagging']}")
        else:
            ckprint("TABLE TAGGING: skipped (requires PDF source + PDF converted)")

        if links_supported(source_path) and links_supported(converted_path):
            if source_path.suffix.lower() == ".pdf":
                src_links = extract_pdf_links(source_path)
            else:
                src_links = extract_html_links(source_path.read_text(encoding="utf-8", errors="ignore"))

            if converted_path.suffix.lower() == ".pdf":
                cnv_links = extract_pdf_links(converted_path)
            else:
                cnv_links = extract_html_links(converted_path.read_text(encoding="utf-8", errors="ignore"))

            missing = sorted(src_links - cnv_links)
            converted_invalid_urls = sorted(
                [u for u in cnv_links if u.startswith(("http://", "https://")) and not is_valid_url(u)])

            result["links"] = {
                "source_count": len(src_links),
                "converted_count": len(cnv_links),
                "missing": missing,
                "converted_invalid_urls": converted_invalid_urls,
            }

            ckprint(f"LINKS source={len(src_links)} converted={len(cnv_links)} missing={len(missing)}")
            if missing:
                print(f"Missing links ({len(missing)}): {missing[:20]}{' ...' if len(missing) > 20 else ''}",
                      flush=True)
            if converted_invalid_urls:
                print(f"Converted invalid URLs ({len(converted_invalid_urls)}): {converted_invalid_urls[:20]}",
                      flush=True)
        else:
            ckprint("LINKS: skipped (only PDF/HTML supported for link extraction)")

        return result

    except Exception as e:
        dump_ui_debug(page, "run_checkpoints_unhandled_exception")
        result["error"] = str(e)
        return result
# =============================================================================
# FAILURE SUMMARY HELPER
# =============================================================================
def checkpoint_failures(result: dict) -> list[str]:

    failures = []

    if result.get("error"):
        failures.append(f"Flow error: {result['error']}")

    if result.get("pages_converted_ok") is False:
        failures.append("Page conversion count mismatch.")

    if result.get("ai_desc_ok") is False:
        failures.append("AI image descriptions did not match expectation.")

    table_tagging = result.get("table_tagging") or {}

    if table_tagging.get("ok") is False:
        failures.append("Table tagging mismatch.")

    links = result.get("links") or {}

    if links.get("missing"):
        failures.append("Converted file is missing source links.")

    if links.get("converted_invalid_urls"):
        failures.append("Converted file contains invalid URLs.")

    return failures
# =============================================================================
# TEST FUNCTIONS
# =============================================================================
@pytest.mark.browser_context_args(storage_state=str(BASE_DIR / "storage_state.json"))
def test_page_count(page):
    "Reads UI page count 'document has X number of pages.'"

def test_ai_descriptions():
    pass
